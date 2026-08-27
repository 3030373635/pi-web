"""Smoke test for the Python document runtime shipped in Linux ARM64 bundles."""

from pathlib import Path
import sys

import defusedxml.ElementTree
from docx import Document
from docxtpl import DocxTemplate
import numpy as np
from openpyxl import Workbook, load_workbook
import pandas as pd
import pdfplumber
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader
import pypdfium2 as pdfium
from reportlab.pdfgen import canvas
import xlsxwriter


def verify_data_packages() -> None:
    """Verify that NumPy, pandas, Pillow, and defusedxml execute real operations."""
    values = np.array([1, 2, 3], dtype=np.int64)
    frame = pd.DataFrame({"value": values})
    assert int(frame["value"].sum()) == 6

    image = Image.new("RGB", (4, 3), color="white")
    assert image.size == (4, 3)

    xml_root = defusedxml.ElementTree.fromstring("<root><value>safe</value></root>")
    assert xml_root.findtext("value") == "safe"


def verify_excel_files(output_root: Path) -> None:
    """Create and reopen Excel files with both engines.

    Args:
        output_root: Directory where smoke-test workbooks are written.
    """
    openpyxl_path = output_root / "openpyxl.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet["A1"] = "Pi Web"
    workbook.save(openpyxl_path)
    assert load_workbook(openpyxl_path).active["A1"].value == "Pi Web"

    xlsxwriter_path = output_root / "xlsxwriter.xlsx"
    generated_workbook = xlsxwriter.Workbook(xlsxwriter_path)
    generated_sheet = generated_workbook.add_worksheet("Data")
    generated_sheet.write("A1", "Pi Web")
    generated_workbook.close()
    assert load_workbook(xlsxwriter_path)["Data"]["A1"].value == "Pi Web"


def verify_word_file(output_root: Path) -> None:
    """Render a DOCX template and verify the generated text.

    Args:
        output_root: Directory where smoke-test Word files are written.
    """
    template_path = output_root / "template.docx"
    rendered_path = output_root / "rendered.docx"
    template_document = Document()
    template_document.add_paragraph("Hello {{ name }}")
    template_document.save(template_path)

    template = DocxTemplate(template_path)
    template.render({"name": "Pi Web"})
    template.save(rendered_path)
    rendered_text = "\n".join(paragraph.text for paragraph in Document(rendered_path).paragraphs)
    assert rendered_text == "Hello Pi Web"


def verify_powerpoint_file(output_root: Path) -> None:
    """Create and reopen a PPTX presentation.

    Args:
        output_root: Directory where the smoke-test presentation is written.
    """
    presentation_path = output_root / "presentation.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Pi Web"
    presentation.save(presentation_path)

    reopened = Presentation(presentation_path)
    assert reopened.slides[0].shapes.title.text == "Pi Web"


def verify_pdf_file(output_root: Path) -> None:
    """Generate, parse, extract, and render a PDF document.

    Args:
        output_root: Directory where the smoke-test PDF is written.
    """
    pdf_path = output_root / "document.pdf"
    pdf_canvas = canvas.Canvas(str(pdf_path))
    pdf_canvas.drawString(72, 720, "Pi Web PDF")
    pdf_canvas.save()

    assert len(PdfReader(pdf_path).pages) == 1
    with pdfplumber.open(pdf_path) as document:
        assert "Pi Web PDF" in (document.pages[0].extract_text() or "")

    # Rendering exercises the bundled native PDFium ARM64 library, not only Python imports.
    pdf_document = pdfium.PdfDocument(pdf_path)
    rendered_image = pdf_document[0].render(scale=1).to_pil()
    assert rendered_image.width > 0 and rendered_image.height > 0


def main(output_directory: str) -> None:
    """Run all document-runtime smoke checks.

    Args:
        output_directory: Directory used for generated verification artifacts.
    """
    output_root = Path(output_directory).resolve()
    output_root.mkdir(parents=True, exist_ok=True)
    verify_data_packages()
    verify_excel_files(output_root)
    verify_word_file(output_root)
    verify_powerpoint_file(output_root)
    verify_pdf_file(output_root)


if __name__ == "__main__":
    if len(sys.argv) != 2:
        raise SystemExit("Usage: python-runtime-smoke.py <output-directory>")
    main(sys.argv[1])
